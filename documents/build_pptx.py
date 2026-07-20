#!/usr/bin/env python3
"""Rebuild 'Short 3 .pptx' → 'Short 4.pptx': corrected metrics, new losses, 10 slides."""
import copy
import lxml.etree as etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# ── Colors ───────────────────────────────────────────────────────────────────
C_DARK_BLUE = RGBColor(0x1F, 0x49, 0x7D)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x1E, 0x1E, 0x1E)
C_GRAY      = RGBColor(0x66, 0x66, 0x66)
C_GREEN     = RGBColor(0x00, 0x6B, 0x00)
C_ROW_EVEN  = RGBColor(0xE8, 0xF0, 0xFA)
C_ROW_ODD   = RGBColor(0xFF, 0xFF, 0xFF)
C_EMPTY_COL = RGBColor(0x99, 0x99, 0x99)

# ── Core helper: clear + rewrite a text frame ─────────────────────────────────

def set_tf(shape, lines, default_size=13):
    """
    Rewrite shape's text frame. Preserves position/size.
    lines: list of str | dict(text, bold, italic, size, color, space_before)
    """
    tf = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody

    # Remove all paragraphs except first, clear first
    paras = txBody.findall(qn('a:p'))
    for p in paras[1:]:
        txBody.remove(p)
    first_p = paras[0]
    for child in list(first_p):
        if child.tag in (qn('a:r'), qn('a:br')):
            first_p.remove(child)

    is_first = True
    for line in lines:
        if isinstance(line, str):
            item = {'text': line}
        else:
            item = line

        if is_first:
            p = first_p
            is_first = False
        else:
            p = etree.SubElement(txBody, qn('a:p'))

        pPr = p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(p, qn('a:pPr'))

        spc = item.get('space_before', 0)
        if spc:
            spcBef = pPr.find(qn('a:spcBef'))
            if spcBef is None:
                spcBef = etree.SubElement(pPr, qn('a:spcBef'))
            spcPts = spcBef.find(qn('a:spcPts'))
            if spcPts is None:
                spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(spc * 100)))

        text  = item.get('text', '')
        bold  = item.get('bold', False)
        italic= item.get('italic', False)
        size  = item.get('size', default_size)
        color = item.get('color', C_BLACK)

        r = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'), {'lang': 'en-US', 'dirty': '0'})
        rPr.set('sz', str(int(size * 100)))
        if bold:    rPr.set('b', '1')
        if italic:  rPr.set('i', '1')
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
        t = etree.SubElement(r, qn('a:t'))
        t.text = text


def add_txbox(slide, l, t, w, h, lines, default_size=13, wrap=True):
    """Add a new text box to a slide."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = wrap
    set_tf(box, lines, default_size)
    return box


def add_table(slide, data, col_widths, left, top, row_height=0.32, font_size=10):
    """Add a styled results table. Bold markers: **text** = green+bold, — = gray."""
    n_rows, n_cols = len(data), len(col_widths)
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(total_w), Inches(n_rows * row_height)
    ).table
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(data):
        for j, raw in enumerate(row):
            cell = tbl.cell(i, j)
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_DARK_BLUE
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_ROW_EVEN
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_ROW_ODD
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            is_bold = raw.startswith('**') and raw.endswith('**')
            is_dash = raw == '—'
            text = raw[2:-2] if is_bold else raw
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = (i == 0 or is_bold)
            if i == 0:
                run.font.color.rgb = C_WHITE
            elif is_bold:
                run.font.color.rgb = C_GREEN
            elif is_dash:
                run.font.color.rgb = C_EMPTY_COL
            else:
                run.font.color.rgb = C_BLACK


def duplicate_slide(prs, src_idx):
    """Duplicate slide[src_idx] and append at end. Returns new slide."""
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sp_tree = new_slide.shapes._spTree
    for ch in list(sp_tree)[2:]:
        sp_tree.remove(ch)
    for ch in list(src.shapes._spTree)[2:]:
        sp_tree.append(copy.deepcopy(ch))
    return new_slide


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    elems = list(lst)
    elem = elems[from_idx]
    lst.remove(elem)
    lst.insert(to_idx, elem)


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation('/home/khang/IR-IID/documents/Short 3 .pptx')

# ── Remove old/stale tables from slides that are being repurposed ─────────────
# Original PPTX had embedded result tables on slides 6, 7, 8, 9 (old data).
# Slides 6→LossII and 7→Metrics&Benchmarks get completely new text content.
# Slide 8→Ablation gets a fresh table; old Table 3 must go.
# Slide 9→IIW has no table in the new design.
SLIDES_TO_PURGE_TABLES = [5, 6, 7, 8]  # 0-indexed: LossII, Metrics, Ablation, IIW
for idx in SLIDES_TO_PURGE_TABLES:
    slide = prs.slides[idx]
    for sh in list(slide.shapes):
        if sh.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
            remove_shape(sh)

# ── S1: Title ─────────────────────────────────────────────────────────────────
s1 = prs.slides[0]
for sh in s1.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Thesis Progress (June 20th)', 'bold': True, 'size': 40, 'color': C_WHITE}])
    elif 'Subtitle' in sh.name:
        set_tf(sh, [
            {'text': 'CARI — Cross-render Albedo-invariant Intrinsics', 'size': 22, 'color': C_WHITE},
            {'text': '(Minh Khang Tran)', 'size': 18, 'color': RGBColor(0xCC, 0xCC, 0xCC)},
        ])

# ── S2: Introduction ──────────────────────────────────────────────────────────
s2 = prs.slides[1]
for sh in s2.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Introduction', 'bold': True, 'size': 32, 'color': C_DARK_BLUE}])
    else:
        set_tf(sh, [
            {'text': 'The problem', 'bold': True, 'size': 15, 'color': C_DARK_BLUE},
            {'text': 'A single photo cannot tell whether a region looks reddish because the material is red or because the light is reddish. One image → infinitely many (Albedo, Shading) pairs explain it.', 'size': 13},
            {'text': '', 'size': 6},
            {'text': 'Measured finding', 'bold': True, 'size': 15, 'color': C_DARK_BLUE, 'space_before': 8},
            {'text': 'State-of-the-art models — including diffusion-based Marigold — still absorb part of the illuminant colour into albedo (measurable via cross-illumination constancy C_mat / C_arap). V17 outperforms Marigold on material constancy (MID C_mat, ARAP C_arap / Cast, MAW intensity) while Marigold leads on structure quality (IIW WHDR, ARAP si-RMSE, MAW ΔE). Both fail to fully separate albedo from illumination; V17 makes the constancy trade-off.', 'size': 13},
            {'text': '', 'size': 6},
            {'text': 'The idea', 'bold': True, 'size': 15, 'color': C_DARK_BLUE, 'space_before': 8},
            {'text': 'Import the disambiguating signal at training time only — from the Multi-Illumination Dataset (MID), which has pixel-aligned multi-light photos of the same scene — and teach a single-image transformer that albedo must stay constant when only the light changes.', 'size': 13},
            {'text': '', 'size': 6},
            {'text': 'Single-image coloured-illumination ambiguity is PARTLY resolvable by cross-render training supervision — without diffusion fine-tuning, without an LLM judge, fully feed-forward.', 'bold': True, 'italic': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 8},
        ], default_size=13)

# ── S3: Our Model (keep images, just refresh text) ────────────────────────────
s3 = prs.slides[2]
for sh in s3.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Our Model', 'bold': True, 'size': 32, 'color': C_DARK_BLUE}])
    elif 'TextBox 5' in sh.name:
        set_tf(sh, [{'text': "What's new is not the architecture — it's how it's trained:", 'bold': True, 'size': 13, 'color': C_DARK_BLUE}])
    elif 'TextBox 13' in sh.name:
        set_tf(sh, [{'text': 'Two images of the same scene under different light are pushed through the same network (Siamese), and two losses connect the two outputs. At test time only one branch runs — single image in, (Albedo, Shading) out.', 'size': 12, 'color': C_GRAY}])

# ── S4: Dataset ───────────────────────────────────────────────────────────────
s4 = prs.slides[3]
for sh in s4.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Dataset', 'bold': True, 'size': 32, 'color': C_DARK_BLUE}])
    else:
        set_tf(sh, [
            {'text': '1.  MIDIntrinsics (MID) — cross-render disambiguation signal', 'bold': True, 'size': 14, 'color': C_DARK_BLUE},
            {'text': '~985 train / 30 test real indoor scenes', 'size': 13},
            {'text': '25 pixel-aligned photos per scene, same camera, 25 different light directions', 'size': 13},
            {'text': 'Shared (pseudo-GT, weak) albedo + per-pixel material-ID map per scene', 'size': 13},
            {'text': 'RAW (un-white-balanced) pairs: MID flash chromaticity varies ~15–25% across directions (measured). This real coloured-illuminant signal is what L_inv must become invariant to. Using synthetic tint augmentation instead measurably HURT (ARAP Cast_RMS 0.035 → 0.078).', 'size': 12, 'italic': True, 'color': C_DARK_BLUE, 'space_before': 3},
            {'text': '', 'size': 5},
            {'text': '2.  Hypersim — synthetic supervised anchor', 'bold': True, 'size': 14, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'Clean per-pixel albedo + shading GT from ray-traced indoor scenes', 'size': 13},
            {'text': 'Prevents the model from collapsing to gray — absolute colour reference throughout training', 'size': 13},
            {'text': '', 'size': 5},
            {'text': 'InteriorVerse excluded (multi-view geometry, not multi-illuminant — no cross-render signal). OpenRooms server unavailable.', 'size': 11, 'color': C_GRAY, 'space_before': 4},
        ], default_size=13)

# ── S5: Loss Design I — Supervised Anchors ────────────────────────────────────
s5 = prs.slides[4]
# Remove equation screenshot images (small images at top of this slide)
for sh in list(s5.shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        remove_shape(sh)
for sh in s5.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Loss Design I — Supervised Anchors', 'bold': True, 'size': 28, 'color': C_DARK_BLUE}])
    else:
        # Resize to fill slide
        sh.left  = Inches(0.67); sh.top    = Inches(1.35)
        sh.width = Inches(12.0); sh.height = Inches(5.80)
        set_tf(sh, [
            {'text': 'Active on BOTH Hypersim (exact GT) and MID (pseudo-GT albedo):', 'bold': True, 'size': 13, 'color': C_DARK_BLUE},
            {'text': '', 'size': 5},
            {'text': 'Albedo  (λ_a = 1.0 + DSSIM 0.2 + MSG 0.5 + chroma L1 0.2)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'L1 data term + DSSIM sharpness + multi-scale gradient (MSG) matching predicted edges to GT + chroma direction L1 penalises hue / saturation error independently of brightness. Chroma term makes desaturation expensive (the MSE-cheap escape).', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'Shading SSI  (λ_s = 1.0)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'Scale-invariant loss on predicted shading (π-domain). Gated to Hypersim (exact GT only). 3000-iter warmup: albedo must anchor first so shading SSI does not push the model to grey-collapse.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'Reconstruction  (λ_recon = 1.0)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'Predicted albedo × shading must recover the input image (A · S ≈ I). Two-sided coupling forces A and S to jointly explain all image structure, closing the "dump residual in R" escape hatch.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'DINOv2 material consistency  (λ = 0.05)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'GT-free. Pixels with similar DINOv2 patch tokens (same material) are penalised if their albedo colours differ. Uses the frozen backbone\'s illumination-invariant features at zero extra inference cost.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'Residual  (λ_r = 0.5 + sparsity 0.02)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'R = (I − A·S)₊ — analytic specular / non-Lambertian remainder. Sparsity prior keeps R near zero outside highlights. On MID, R_star ≡ 0 by construction (closes the "dump colour into R" hatch).', 'size': 12},
        ], default_size=12)

# ── S6: Loss Design II — CARI  (repurpose sparse Metrics slide) ───────────────
s6 = prs.slides[5]
for sh in s6.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Loss Design II — CARI (the Contribution)', 'bold': True, 'size': 28, 'color': C_DARK_BLUE}])
    else:
        sh.left  = Inches(0.67); sh.top    = Inches(1.35)
        sh.width = Inches(12.0); sh.height = Inches(5.80)
        set_tf(sh, [
            {'text': 'Active from Phase 3 — requires paired (rgb, rgb₂): SAME scene, DIFFERENT light direction.', 'bold': True, 'size': 13, 'color': C_DARK_BLUE},
            {'text': '', 'size': 5},
            {'text': 'L_inv — Cross-render albedo invariance  (λ = 0.5)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'Same surface, different light → albedo must be IDENTICAL.  Any difference is the model wrongly absorbing the illuminant into albedo. This is the term that directly attacks the coloured-cast problem.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'L_explain — Shading explains the change  (λ = 0.25)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'Whatever changed between rgb and rgb₂ is ONLY the light — forces the entire inter-frame intensity ratio into shading, not albedo. Log-domain ratio, robust to scale. Prevents the complement failure: CARI keeps albedo constant but shading absorbs the residual.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'L_shade_sign  (λ = 0.05)  — data-free, all phases', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'A · S must not exceed I anywhere — prevents the model from brightening albedo beyond the observed image.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'Why RAW MID pairs?', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 6},
            {'text': 'MID flash directions are NOT white — illuminant chromaticity MEASURED to vary ~15–25% across 25 directions. Using un-white-balanced (raw) pairs forces L_inv to handle REAL coloured illuminants rather than synthetic white-only perturbations. Chromatic augmentation U[0.6, 1.4] measurably HURT (ARAP Cast_RMS 0.035 → 0.078 at 40k).', 'size': 12, 'italic': True, 'color': C_DARK_BLUE},
        ], default_size=12)

# ── S7: Metrics & Benchmarks  (repurpose sparse Benchmarks slide) ─────────────
s7 = prs.slides[6]
for sh in s7.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Metrics & Benchmarks', 'bold': True, 'size': 32, 'color': C_DARK_BLUE}])
    else:
        sh.left  = Inches(0.67); sh.top    = Inches(1.35)
        sh.width = Inches(12.0); sh.height = Inches(5.80)
        set_tf(sh, [
            {'text': '4 benchmarks — all metrics LOWER = better', 'bold': True, 'size': 14, 'color': C_DARK_BLUE},
            {'text': '', 'size': 5},
            {'text': 'MID constancy  (30 test scenes, real indoor)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 5},
            {'text': 'C_mat: std-dev of albedo chroma across 25 lights per material region → coloured-cast invariance (the thesis axis).', 'size': 12},
            {'text': 'Cast_RMS: global RG / BG channel ratio drift across relightings.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'ARAP constancy  (37 scene groups, synthetic, out-of-domain)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 5},
            {'text': 'C_arap / Cast_RMS: same constancy axes with EXACT synthetic GT. Model never trains on ARAP — fair cross-domain check.', 'size': 12},
            {'text': 'si-RMSE: scale-invariant RMSE vs exact GT albedo → STRUCTURE accuracy (texture sharpness; guard metric against gray-collapse).', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'MAW  (874 images, real diverse / wild)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 5},
            {'text': 'ΔE: CIE chromaticity error per material patch.  Intensity SI-MSE×100: scale-invariant brightness error.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'IIW  (200 images, real held-out)', 'bold': True, 'size': 13, 'color': C_DARK_BLUE, 'space_before': 5},
            {'text': 'WHDR: Weighted Human Disagreement Rate — standard pairwise reflectance ordering accuracy.', 'size': 12},
            {'text': '', 'size': 5},
            {'text': 'Accuracy (si-RMSE) and constancy (C_mat / C_arap) are always reported together: a model can trivially "win" constancy by predicting flat gray — accuracy rules that out.', 'size': 11, 'italic': True, 'color': C_GRAY, 'space_before': 4},
            {'text': 'Why both MID and ARAP: MID is real but its albedo is only pseudo-GT (proves the mechanism); ARAP is synthetic with exact GT and is out-of-domain (the harder cross-check).', 'size': 11, 'italic': True, 'color': C_GRAY},
        ], default_size=12)

# ── NEW S8: Results table  (duplicate slide 7 as base, wipe, add content) ─────
new_res = duplicate_slide(prs, 6)   # appended at index 9 (after slides 0-8)

# Wipe all shapes from the duplicate
for sh in list(new_res.shapes):
    remove_shape(sh)

# Title
add_txbox(new_res, 0.67, 0.28, 12.0, 0.70,
          [{'text': 'Results — V17 vs. Marigold', 'bold': True, 'size': 32, 'color': C_DARK_BLUE}])

# Subtitle note
add_txbox(new_res, 0.67, 0.95, 12.0, 0.40,
          [{'text': 'All metrics LOWER = better.  ARAP = indoor split (23 groups).  Marigold output linearised to linear space before scoring (output-space fix).  Bold green = best per column.', 'size': 10, 'italic': True, 'color': C_GRAY}])

# Results table
table_data = [
    ['Model',                  'MID\nC_mat ↓', 'MID\nCast ↓', 'ARAP\nC_arap ↓', 'ARAP\nCast ↓', 'ARAP\nsi-RMSE ↓', 'MAW\nΔE ↓', 'MAW\nInt×100 ↓', 'IIW\nWHDR ↓'],
    ['V17  (19k, no CARI)',     '0.150',         '0.131',        '**0.151**',        '**0.063**',    '0.332',            '4.779',     '0.427',           '**0.290**'],
    ['V17  (50k, full CARI)',   '**0.107**',      '0.136',        '0.166',            '0.071',        '0.362',            '4.023',     '**0.402**',        '0.303'],
    ['Marigold-app',            '0.140',          '**0.085**',    '0.170',            '0.125',        '**0.283**',        '**3.625**', '0.461',            '**0.130**'],
    ['Marigold-light',          '0.276',          '0.126',        '0.292',            '—',            '0.576',            '4.047',     '0.422',            '0.222'],
]
col_w = [2.05, 1.10, 1.10, 1.20, 1.15, 1.30, 1.05, 1.35, 1.10]
add_table(new_res, table_data, col_w, left=0.30, top=1.42, row_height=0.72, font_size=10.5)

# Summary note
add_txbox(new_res, 0.67, 5.50, 12.0, 0.80, [
    {'text': 'V17 wins constancy axes (MID C_mat, ARAP C_arap / Cast, MAW Intensity) — the thesis contribution.  Marigold-app wins structure axes (IIW WHDR, ARAP si-RMSE, MAW ΔE).', 'bold': True, 'size': 11, 'color': C_DARK_BLUE},
    {'text': 'Marigold-light is trained on Hypersim only (no cross-render), like our 19k baseline. Its ARAP Cast is excluded (Cast_RMS ≈ 10.85 — sign error, broken metric for that model).', 'size': 10, 'italic': True, 'color': C_GRAY},
])

# ── S8 (currently index 7): Ablation Study — keep viz images, refresh text ────
s_abl = prs.slides[7]
for sh in s_abl.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Ablation Study — MID Constancy', 'bold': True, 'size': 30, 'color': C_DARK_BLUE}])
    elif 'TextBox 5' in sh.name:
        sh.left  = Inches(0.30); sh.top    = Inches(1.28)
        sh.width = Inches(12.80); sh.height = Inches(0.35)
        set_tf(sh, [{'text': 'MID C_mat: Row 1 (19k) = 0.150 → Row 4 (50k, full CARI) = 0.107  (−28%).  Rows 2 & 3 results pending.', 'bold': True, 'size': 11, 'color': C_DARK_BLUE}])

# Add ablation table between title and the visualization images
# Images start at y ≈ 2.86" — table fits nicely at top=1.62"
abl_data = [
    ['Row', 'Config',                     'MID C_mat ↓', 'MID Cast ↓', 'ARAP C_arap ↓', 'ARAP Cast ↓', 'IIW WHDR ↓'],
    ['1',   'Supervised only  (19k)',      '0.150',        '0.131',       '0.151',           '0.063',        '0.290'],
    ['2',   '+ L_inv + L_explain  (30k)',  '—',            '—',           '—',               '—',            '—'],
    ['3',   '+ albedo RGB-skip  (30k)',    '—',            '—',           '—',               '—',            '—'],
    ['4',   'Full CARI  (50k)',            '**0.107**',    '0.136',       '0.166',           '0.071',        '0.303'],
]
abl_col_w = [0.45, 2.80, 1.15, 1.10, 1.25, 1.15, 1.10]
add_table(s_abl, abl_data, abl_col_w, left=0.30, top=1.63, row_height=0.24, font_size=9.5)

# Row labels left of the three visualization image strips
for label, y_in in [
    ('Row 1  (19k)',  2.72),
    ('Row 2  (30k)',  4.05),
    ('Row 4  (50k)',  5.50),
]:
    add_txbox(s_abl, 0.00, y_in, 2.60, 0.22,
              [{'text': label, 'size': 8.5, 'bold': True, 'color': C_DARK_BLUE}])

# ── S9 (currently index 8): IIW — update text ─────────────────────────────────
s_iiw = prs.slides[8]
for sh in s_iiw.shapes:
    if not sh.has_text_frame: continue
    if 'Title' in sh.name:
        set_tf(sh, [{'text': 'Ablation Study — IIW (no-regression check)', 'bold': True, 'size': 28, 'color': C_DARK_BLUE}])
    elif 'TextBox 11' in sh.name:
        sh.left  = Inches(0.50); sh.top    = Inches(1.38)
        sh.width = Inches(12.50); sh.height = Inches(0.65)
        set_tf(sh, [
            {'text': 'IIW WHDR (200 images): Row 1 (19k) = 29.0%  →  Row 4 (50k) = 30.3%  |  Marigold-app = 13.0%', 'bold': True, 'size': 12, 'color': C_DARK_BLUE},
            {'text': 'Slight regression with full CARI training. Marigold-app substantially outperforms on IIW (structure quality axis — diffusion edge sharpness). CARI did not catastrophically hurt IIW.', 'size': 11, 'color': C_GRAY},
        ])

# ── Reorder: move the new Results slide (now at index 9) to index 7 ───────────
# Desired final order: 0(Title) 1(Intro) 2(Model) 3(Dataset) 4(LossI) 5(LossII) 6(Metrics) 7(Results) 8(Ablation) 9(IIW)
move_slide(prs, 9, 7)

# ── Save ──────────────────────────────────────────────────────────────────────
out = '/home/khang/IR-IID/documents/Short 4.pptx'
prs.save(out)
print(f'Saved → {out}')
print(f'Slides: {len(prs.slides)}')
